import os
import streamlit as st
from dotenv import load_dotenv
from pptx import Presentation
from pptx.util import Inches
from agents.pitch_deck_agent import PitchDeckAgent

# Φόρτωση μεταβλητών περιβάλλοντος
load_dotenv()

# Streamlit settings
st.set_page_config(page_title="Startup Finance Copilot", layout="wide")
st.title("🚀 Startup Finance Copilot")
st.subheader("Δημιουργία Pitch Deck με Οικονομικές Προβλέψεις")

# Ορισμός διαδρομής αποθήκευσης
DATA_DIR = "data"
os.makedirs(DATA_DIR, exist_ok=True)

# Φόρτωμα αρχείων
uploaded_txt = st.file_uploader("📄 Ανέβασε αρχείο με πληροφορίες για τη startup (.txt)", type=["txt"])
uploaded_excel = st.file_uploader("📊 Ανέβασε οικονομικά δεδομένα (.xlsx)", type=["xlsx"])

if uploaded_txt and uploaded_excel:
    with st.spinner("🔄 Επεξεργασία αρχείων..."):
        # Ονόματα για αποθήκευση
        txt_path = os.path.join(DATA_DIR, "startup_info.txt")
        excel_path = os.path.join(DATA_DIR, "financials.xlsx")

        # Αποθήκευση αρχείων
        with open(txt_path, "wb") as f:
            f.write(uploaded_txt.read())
        with open(excel_path, "wb") as f:
            f.write(uploaded_excel.read())

        # Δημιουργία pitch deck
        agent = PitchDeckAgent()
        result = agent.generate_pitch_deck(txt_path=txt_path, excel_path=excel_path)
        print("Αποτέλεσμα από pitch deck agent:")
        print(result)

    st.success("✅ Το pitch deck δημιουργήθηκε με επιτυχία!")
    st.markdown("### 🧾 Αποτέλεσμα:")
    st.write(result)

    # Δημιουργία παρουσίασης
    prs = Presentation()

    slides = result.split("\n\n")

    for slide_text in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        body = slide.shapes.placeholders[1]
        title.text = slide_text.split("\n")[0]
        body.text = "\n".join(slide_text.split("\n")[1:])

    prs.save("pitch_deck.pptx")

    with open("pitch_deck.pptx", "rb") as f:
        st.download_button("Κατέβασε το Pitch Deck σε PowerPoint", f, file_name="pitch_deck.pptx")



else:
    st.info("📥 Παρακαλώ ανέβασε και τα δύο αρχεία για να συνεχίσεις.")


# streamlit run startup-finance-copilot/app.py